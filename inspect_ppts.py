# -*- coding: utf-8 -*-
"""Extract all embedded images from PPT files + dump slide text per file."""
import sys, io, os
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
import os, glob

PPT_DIR = r"C:/Users/unive/OneDrive/111 유튜버/학생 흉부 사진/20260519 학생"
DUMP_DIR = r"C:/Users/unive/OneDrive/111 Claude/병동 가이드 라인/cxr-grading-app/_ppt_dump"
os.makedirs(DUMP_DIR, exist_ok=True)

ppts = sorted(glob.glob(os.path.join(PPT_DIR, "*.pptx")))
print(f"Found {len(ppts)} PPT files\n")

for pi, ppt_path in enumerate(ppts):
    name = os.path.basename(ppt_path).replace(".pptx", "")
    safe = name.replace(" ", "_").replace("/", "_")[:30]
    print(f"=" * 70)
    print(f"[{pi:02d}] {name}")
    print(f"=" * 70)
    try:
        prs = Presentation(ppt_path)
    except Exception as e:
        print(f"  ERROR: {e}")
        continue
    for si, slide in enumerate(prs.slides):
        # Text from slide
        texts = []
        img_count_on_slide = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    line = "".join(r.text for r in p.runs).strip()
                    if line:
                        texts.append(line)
            if shape.shape_type == 13:  # picture
                img_count_on_slide += 1
                blob = shape.image.blob
                ext = shape.image.ext  # png/jpeg/etc.
                fn = f"{pi:02d}_{safe}_s{si:02d}_i{img_count_on_slide}.{ext}"
                with open(os.path.join(DUMP_DIR, fn), "wb") as f:
                    f.write(blob)
        snippet = " | ".join(texts)[:200]
        if texts or img_count_on_slide:
            print(f"  slide {si+1:02d}: imgs={img_count_on_slide} | text: {snippet}")
print(f"\nImages dumped to: {DUMP_DIR}")
