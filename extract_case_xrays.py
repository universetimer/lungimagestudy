# -*- coding: utf-8 -*-
"""
For each unique student PPT:
  - Find each case's "comparison" slide (after a title slide, has 2 images,
    text mentions "Normal chest PA")
  - Extract the 2nd image (= current/abnormal X-ray)
  - Also extract findings text from the next slide (12-item reading)
  - Save image + write JSON metadata
"""
import sys, io, os, json, shutil
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPT_DIR = r"C:/Users/unive/OneDrive/111 유튜버/학생 흉부 사진/20260519 학생"
OUT_IMG_DIR = r"C:/Users/unive/OneDrive/111 Claude/병동 가이드 라인/cxr-grading-app/_case_xrays"
OUT_JSON = r"C:/Users/unive/OneDrive/111 Claude/병동 가이드 라인/cxr-grading-app/_extracted_cases.json"

# Unique student PPTs (skip duplicates)
PPTS = [
    ("hong",  "CXR (1) 홍예림.pptx"),
    ("lee",   "CXR (2) 이현수.pptx"),
    ("seo",   "CXR (3)서동성.pptx"),
    ("yu",    "CXR 발표_A반 19조 57번 202300060 유병진.pptx"),
    ("mok",   "CXR 판독_A반 20조 60번 202300028 목승민.pptx"),
    ("choi",  "CXR 판독_A반 20조 61번 202300095 최시연.pptx"),
]

if os.path.exists(OUT_IMG_DIR):
    shutil.rmtree(OUT_IMG_DIR)
os.makedirs(OUT_IMG_DIR)

def slide_text(slide):
    out = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                t = "".join(r.text for r in p.runs).strip()
                if t: out.append(t)
    return out

def slide_pictures(slide):
    pics = []
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics.append(shape)
    return pics

def is_title_slide(texts):
    # Title slide for a case = short text with patient initials + CC
    # e.g., "김O현 (77/M)" + "CC. Dyspnea, Cough, Sputum"
    if len(texts) > 4: return False
    joined = " ".join(texts)
    return ("CC." in joined or "CC:" in joined) and "/" in joined

def is_comparison_slide(texts, n_pics):
    # Has 2+ images and mentions "Normal chest PA" / "Normal Chest PA"
    if n_pics < 2: return False
    joined = " ".join(texts).lower()
    return "normal chest pa" in joined or "normal chest" in joined

def is_findings_slide(texts):
    # Has the 12-item template — typically starts with "Chest PA" or "Chest AP"
    # AND mentions "Lung volume"
    joined = "\n".join(texts).lower()
    return "lung volume" in joined and ("trachea" in joined or "tracheal" in joined)

def parse_demographics(text):
    # extract age/sex e.g. "(77/M)" → {"age": 77, "sex": "M"}
    import re
    m = re.search(r"\((\d{1,3})\s*/\s*([MFWmfw])\)", text)
    if m:
        return {"age": int(m.group(1)), "sex": m.group(2).upper()}
    return None

def parse_cc(texts):
    for t in texts:
        if "CC." in t or "CC:" in t:
            return t.split("CC.")[-1].split("CC:")[-1].strip().strip("-").strip()
    return ""

results = []

for student_key, fname in PPTS:
    ppt_path = os.path.join(PPT_DIR, fname)
    print(f"\n=== {student_key}: {fname} ===")
    prs = Presentation(ppt_path)
    slides = list(prs.slides)

    case_idx = 0
    i = 0
    while i < len(slides):
        texts = slide_text(slides[i])
        pics = slide_pictures(slides[i])
        if is_title_slide(texts) and len(pics) == 0:
            case_idx += 1
            demo = None
            cc = ""
            for t in texts:
                d = parse_demographics(t)
                if d: demo = d
            cc = parse_cc(texts)

            # Look at next 1-3 slides for comparison + findings
            cmp_slide_idx = None
            findings_text = None
            for j in range(i+1, min(i+4, len(slides))):
                jt = slide_text(slides[j])
                jp = slide_pictures(slides[j])
                if cmp_slide_idx is None and is_comparison_slide(jt, len(jp)):
                    cmp_slide_idx = j
                if is_findings_slide(jt):
                    findings_text = "\n".join(jt)

            if cmp_slide_idx is not None:
                cmp_pics = slide_pictures(slides[cmp_slide_idx])
                # 2nd picture = current (abnormal). Order by left position.
                cmp_pics.sort(key=lambda s: s.left if s.left is not None else 0)
                target = cmp_pics[-1]  # rightmost = newest / current
                # Save image
                blob = target.image.blob
                ext = target.image.ext
                case_id = f"{student_key}_{case_idx}"
                out_path = os.path.join(OUT_IMG_DIR, f"{case_id}.{ext}")
                with open(out_path, "wb") as f:
                    f.write(blob)
                results.append({
                    "case_id": case_id,
                    "student": student_key,
                    "case_in_student": case_idx,
                    "age": demo["age"] if demo else None,
                    "sex": demo["sex"] if demo else None,
                    "cc": cc,
                    "image_file": os.path.basename(out_path),
                    "image_size_bytes": len(blob),
                    "findings_text": findings_text,
                    "comparison_slide": cmp_slide_idx + 1,
                })
                print(f"  case {case_idx}: ({demo['age'] if demo else '?'}/{demo['sex'] if demo else '?'}) CC: {cc[:40]} -> {os.path.basename(out_path)} ({len(blob)//1024}KB)")
            else:
                print(f"  case {case_idx}: no comparison slide found")
        i += 1

with open(OUT_JSON, "w", encoding="utf-8") as f:
    json.dump(results, f, ensure_ascii=False, indent=2)
print(f"\nTotal: {len(results)} cases. Metadata: {OUT_JSON}")
